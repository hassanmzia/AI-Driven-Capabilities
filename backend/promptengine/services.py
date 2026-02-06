"""
Core prompt engineering services - implements all 5 notebook features
plus custom prompt execution.
"""
import io
import re
import time
import json
import logging
from django.conf import settings

logger = logging.getLogger(__name__)

# System prompts derived from the notebook
SYSTEM_PROMPTS = {
    'feedback_analysis': (
        "You are an assistant that supports a healthcare provider in analyzing patient reviews.\n\n"
        "Your goal is to extract key information from the user data provided below, including the "
        "patient's name, the doctor mentioned in the review, the review rating, a brief description "
        "of the review, and whether the patient expressed satisfaction with their appointment. "
        "Go through the user's feedback step by step, and generate a structured output for further "
        "analysis by the healthcare provider in the below format:\n\n"
        "{\n"
        '    "patient_name": "<extract the patient\'s first and last name from the corpus>",\n'
        '    "consulting_doctor": "<extract the doctor\'s first and last name and credentials from the corpus>",\n'
        '    "review_rating": "<this has to be a number out of 5 points - if you cannot find a rating, output NULL>",\n'
        '    "review_description": "<summarize the review at most in 50 words>",\n'
        '    "satisfaction": "<this has to be a TRUE or FALSE value - arrive at this conclusion using your own judgment>",\n'
        '    "issue_tags": "<in the case of a negative review or dissatisfaction, add tags which specify the area of dissatisfaction>"\n'
        "}"
    ),

    'meeting_summarizer': (
        "Your task is to summarize and document meeting transcripts.\n\n"
        "Carefully read through the user input and provide an output with the below sections:\n\n"
        "1. Date of the Meeting\n"
        "2. A summary of the overall objective\n"
        "3. The list of participants and their roles in the organization\n"
        "4. Crisp discussion points\n"
        "5. Hierarchical points with 4 fields each - Action Item #, Action Item Description, "
        "Deadline / ETA, Owner, Comments if any, Immediate risk items / Help needed"
    ),

    'quiz_generator': (
        "You are an assistant to a Corporate Trainer. Your task is to generate multiple choice "
        "questions from a provided text.\n\n"
        "Follow these instructions strictly when you create the quiz:\n\n"
        "1. The questions must come only from the provided text\n"
        "2. A question must have only one correct answer\n"
        "3. Each question must have exactly 4 options (A, B, C, D)\n"
        "4. Clearly indicate the correct answer for each question\n"
        "5. Create the requested difficulty mix\n"
        "6. Format output as JSON array with fields: question, options (array), correct_answer, difficulty"
    ),

    'slide_script': (
        "You are a professional presentation consultant. Generate a concise and compelling "
        "slide script based on the user's requirements.\n\n"
        "For each slide provide:\n"
        "- Slide number\n"
        "- Title\n"
        "- Key bullet points (3-5)\n"
        "- Speaker notes (2-3 sentences for verbal delivery)\n\n"
        "Format output as JSON array with fields: slide_number, title, bullets (array), speaker_notes"
    ),

    'complaint_response': (
        "You are {agent_name}, a Customer Support Assistant for {company_name}.\n"
        "You will be provided with a user complaint.\n"
        "Your task is to generate a response to the user complaint. Please process the user's complaint, "
        "identify its sentiment, and respond appropriately, always maintaining a polite and helpful tone. "
        "If possible, provide relevant information or troubleshooting steps.\n\n"
        "If the sentiment is negative or critical, assure the user that their concern is taken seriously, "
        "and a representative will address it as soon as possible.\n"
        "If the sentiment is positive or neutral, acknowledge the issue and offer any immediate assistance "
        "or guidance available.\n\n"
        "Also output a JSON block at the end with: sentiment, urgency (low/medium/high), category"
    ),
}


def get_llm(model_name='gpt-4o-mini', temperature=0.0, max_tokens=1024):
    """Create and return an LLM instance based on model name."""
    try:
        if 'claude' in model_name.lower() or 'anthropic' in model_name.lower():
            from langchain_anthropic import ChatAnthropic
            return ChatAnthropic(
                model=model_name,
                temperature=temperature,
                max_tokens=max_tokens,
                api_key=settings.ANTHROPIC_API_KEY,
            )
        else:
            from langchain_openai import ChatOpenAI
            return ChatOpenAI(
                model=model_name,
                temperature=temperature,
                max_tokens=max_tokens,
                api_key=settings.OPENAI_API_KEY,
            )
    except Exception as e:
        logger.warning(f"Failed to create LLM ({model_name}): {e}. Using mock response.")
        return None


def execute_prompt(system_prompt, user_prompt, model='gpt-4o-mini', temperature=0.0, max_tokens=1024):
    """Execute a prompt and return the result with metadata."""
    start_time = time.time()

    llm = get_llm(model, temperature, max_tokens)

    if llm is None:
        # Return a mock response when no API key is configured
        elapsed = int((time.time() - start_time) * 1000)
        return {
            'output': (
                f"[Demo Mode - No API key configured]\n\n"
                f"This is a simulated response for the given prompt. "
                f"Configure OPENAI_API_KEY or ANTHROPIC_API_KEY in your .env file "
                f"to get real AI responses.\n\n"
                f"System prompt: {system_prompt[:100]}...\n"
                f"User input length: {len(user_prompt)} chars"
            ),
            'tokens_input': len(user_prompt.split()),
            'tokens_output': 50,
            'cost_estimate': 0.0,
            'latency_ms': elapsed,
            'model': model,
        }

    try:
        messages = []
        if system_prompt:
            messages.append(('system', system_prompt))
        messages.append(('human', user_prompt))

        response = llm.invoke(messages)
        elapsed = int((time.time() - start_time) * 1000)

        tokens_in = getattr(response, 'usage_metadata', {}).get('input_tokens', len(user_prompt.split()))
        tokens_out = getattr(response, 'usage_metadata', {}).get('output_tokens', len(response.content.split()))

        return {
            'output': response.content,
            'tokens_input': tokens_in,
            'tokens_output': tokens_out,
            'cost_estimate': _estimate_cost(model, tokens_in, tokens_out),
            'latency_ms': elapsed,
            'model': model,
        }
    except Exception as e:
        elapsed = int((time.time() - start_time) * 1000)
        logger.error(f"LLM execution failed: {e}")
        return {
            'output': '',
            'error': str(e),
            'tokens_input': 0,
            'tokens_output': 0,
            'cost_estimate': 0.0,
            'latency_ms': elapsed,
            'model': model,
        }


def _estimate_cost(model, tokens_in, tokens_out):
    """Rough cost estimation based on model pricing."""
    pricing = {
        'gpt-4o-mini': (0.00015, 0.0006),
        'gpt-4o': (0.005, 0.015),
        'gpt-4': (0.03, 0.06),
        'claude-3-5-sonnet-20241022': (0.003, 0.015),
        'claude-3-haiku-20240307': (0.00025, 0.00125),
    }
    rates = pricing.get(model, (0.001, 0.002))
    return round((tokens_in / 1000 * rates[0]) + (tokens_out / 1000 * rates[1]), 6)


# --- Feature implementations ---

def analyze_feedback(review_text, model='gpt-4o-mini', temperature=0.0):
    system_prompt = SYSTEM_PROMPTS['feedback_analysis']
    user_prompt = f"---\nHere is the data:\n{review_text}"
    return execute_prompt(system_prompt, user_prompt, model, temperature)


def summarize_meeting(transcript, model='gpt-4o-mini', temperature=0.0):
    system_prompt = SYSTEM_PROMPTS['meeting_summarizer']
    user_prompt = f"---\nBelow is the transcript:\n{transcript}"
    return execute_prompt(system_prompt, user_prompt, model, temperature)


def generate_quiz(content, num_questions=5, difficulty_mix='2 easy, 2 intermediate, 1 hard',
                  model='gpt-4o-mini', temperature=0.0):
    system_prompt = SYSTEM_PROMPTS['quiz_generator']
    user_prompt = (
        f"Generate {num_questions} questions with the following difficulty distribution: {difficulty_mix}\n"
        f"---\nBelow is the information on which you have to generate the Quiz:\n{content}"
    )
    return execute_prompt(system_prompt, user_prompt, model, temperature, max_tokens=2048)


def generate_slide_script(topic, num_slides=3, style='professional',
                          model='gpt-4o-mini', temperature=0.7):
    system_prompt = SYSTEM_PROMPTS['slide_script']
    user_prompt = (
        f"Generate a {style} slide script on the topic: {topic}\n"
        f"Create exactly {num_slides} slides.\n"
        f"Please ensure that each slide contains enough information to provide a clear overview."
    )
    return execute_prompt(system_prompt, user_prompt, model, temperature)


def generate_complaint_response(complaint, company_name='Our Company', agent_name='Support Agent',
                                 model='gpt-4o-mini', temperature=0.3):
    system_prompt = SYSTEM_PROMPTS['complaint_response'].format(
        agent_name=agent_name, company_name=company_name
    )
    user_prompt = f"User complaint:\n{complaint}"
    return execute_prompt(system_prompt, user_prompt, model, temperature)


def execute_custom_prompt(system_prompt='', user_prompt='', model='gpt-4o-mini',
                          temperature=0.7, max_tokens=1024):
    return execute_prompt(system_prompt, user_prompt, model, temperature, max_tokens)


# --- Advanced Prompt Engineering System Prompts ---

ADVANCED_PROMPTS = {
    'prompt_grader': (
        "You are a senior prompt engineering expert. Analyze the given prompt and return a JSON object with:\n"
        "{\n"
        '  "scores": {\n'
        '    "clarity": <1-10>,\n'
        '    "specificity": <1-10>,\n'
        '    "constraints": <1-10>,\n'
        '    "output_format": <1-10>,\n'
        '    "examples": <1-10>,\n'
        '    "safety": <1-10>\n'
        "  },\n"
        '  "overall_grade": "<A/B/C/D/F>",\n'
        '  "gaps": [\n'
        '    {"category": "<dimension>", "issue": "<what is missing or weak>", "suggestion": "<specific fix>"}\n'
        "  ],\n"
        '  "summary": "<2-3 sentence assessment>"\n'
        "}\n\n"
        "Scoring rubric:\n"
        "- Clarity: Is the intent unambiguous? Are instructions clear?\n"
        "- Specificity: Are requirements detailed enough to produce consistent output?\n"
        "- Constraints: Are boundaries, limits, and edge cases defined?\n"
        "- Output Format: Is the expected output structure specified (JSON, sections, etc.)?\n"
        "- Examples: Are few-shot examples or expected I/O provided?\n"
        "- Safety: Are there guardrails against harmful, biased, or off-topic output?\n\n"
        "Return ONLY valid JSON, no markdown or extra text."
    ),

    'prompt_optimizer': (
        "You are a prompt optimization expert. Given the original prompt and its analysis, "
        "produce an improved version that addresses all identified gaps.\n\n"
        "Return a JSON object with:\n"
        "{\n"
        '  "improved_prompt": "<the rewritten prompt preserving original intent but fixing all issues>",\n'
        '  "advanced_prompt": "<a version using few-shot examples, chain-of-thought, and role+rubric patterns>",\n'
        '  "changes": [\n'
        '    {"what": "<what changed>", "why": "<reason for the change>"}\n'
        "  ]\n"
        "}\n\n"
        "Guidelines:\n"
        "- Preserve the original intent completely\n"
        "- Add missing role/persona if absent\n"
        "- Add explicit output format constraints\n"
        "- Add safety guardrails if missing\n"
        "- The advanced version should include few-shot examples and chain-of-thought instructions\n"
        "- Return ONLY valid JSON."
    ),

    'ab_judge': (
        "You are an impartial evaluator comparing two AI outputs for the same task.\n\n"
        "Score each output on these dimensions (1-10):\n"
        "1. Accuracy: Correctness and factual reliability\n"
        "2. Completeness: Coverage of all required elements\n"
        "3. Clarity: Readability and organization\n"
        "4. Relevance: Stays on-topic, no extraneous content\n"
        "5. Format: Follows requested output structure\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "output_a_scores": {"accuracy": <n>, "completeness": <n>, "clarity": <n>, "relevance": <n>, "format": <n>},\n'
        '  "output_b_scores": {"accuracy": <n>, "completeness": <n>, "clarity": <n>, "relevance": <n>, "format": <n>},\n'
        '  "winner": "<A or B or TIE>",\n'
        '  "reasoning": "<brief explanation of why one is better>"\n'
        "}\n"
        "Return ONLY valid JSON."
    ),

    'schema_enforcer': (
        "You are a precise data extraction assistant. You MUST return output that strictly "
        "conforms to the JSON schema provided. Follow these rules:\n"
        "1. Return ONLY valid JSON matching the schema exactly\n"
        "2. Use the exact field names from the schema\n"
        "3. Respect all type constraints (string, number, boolean, array)\n"
        "4. Include all required fields\n"
        "5. If a field value cannot be determined, use null\n"
        "6. Do NOT include any text outside the JSON object\n"
    ),

    'schema_validator_feedback': (
        "The previous output failed JSON schema validation with these errors:\n"
        "{errors}\n\n"
        "Please fix the output to match the required schema. Return ONLY valid JSON."
    ),

    'self_correct_critic': (
        "You are a critical reviewer. Analyze the following AI-generated response and identify issues.\n\n"
        "Evaluate on:\n"
        "1. Factual accuracy - any incorrect claims?\n"
        "2. Completeness - anything missing from the task requirements?\n"
        "3. Logic - any reasoning flaws or contradictions?\n"
        "4. Clarity - any confusing or poorly explained parts?\n"
        "5. Format - does it match the requested output format?\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "quality_score": <1-10>,\n'
        '  "issues": [\n'
        '    {"type": "<accuracy|completeness|logic|clarity|format>", "description": "<specific issue>", "severity": "<high|medium|low>"}\n'
        "  ],\n"
        '  "passes_threshold": <true if quality_score >= 7>,\n'
        '  "revision_instructions": "<specific instructions for improving the response>"\n'
        "}\n"
        "Return ONLY valid JSON."
    ),

    'self_correct_reviser': (
        "You are a revision expert. Given the original response and the critic's feedback, "
        "produce an improved version that addresses ALL identified issues.\n\n"
        "Maintain the original intent and format while fixing every issue noted by the critic. "
        "Do not mention that this is a revision - just output the improved response directly."
    ),

    'quality_gate_safety': (
        "You are a content safety and policy reviewer. Analyze the following content for:\n"
        "1. Harmful or offensive content\n"
        "2. Bias or discrimination\n"
        "3. Factual inaccuracies or misleading claims\n"
        "4. Privacy violations (PII exposure)\n"
        "5. Completeness - does it address all requirements?\n"
        "6. Professional tone and appropriateness\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "passed": <true/false>,\n'
        '  "checks": [\n'
        '    {"check": "<check name>", "status": "<pass|fail|warning>", "details": "<explanation>"}\n'
        "  ],\n"
        '  "revision_needed": "<specific fixes required if failed, or empty string>"\n'
        "}\n"
        "Return ONLY valid JSON."
    ),

    'quality_gate_reviser': (
        "You are a content editor. The following content failed a quality review. "
        "Revise it to address ALL identified issues while preserving the original meaning and intent.\n\n"
        "Issues to fix:\n{issues}\n\n"
        "Output ONLY the revised content, no meta-commentary."
    ),

    'decomposition_planner': (
        "You are a task decomposition expert. Break the following complex task into "
        "clear, sequential sub-tasks that can each be handled independently.\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "task_analysis": "<brief analysis of the overall task>",\n'
        '  "sub_tasks": [\n'
        '    {"step": <number>, "title": "<short title>", "instruction": "<detailed instruction for this sub-task>", "depends_on": "<what input this step needs>"}\n'
        "  ]\n"
        "}\n"
        "Create 3-6 sub-tasks. Return ONLY valid JSON."
    ),

    'decomposition_integrator': (
        "You are a synthesis expert. Multiple sub-tasks have been completed for a larger task. "
        "Integrate all the sub-task outputs into a single, cohesive final result.\n\n"
        "Ensure the final output:\n"
        "- Flows naturally as a unified document\n"
        "- Removes any redundancy between sub-tasks\n"
        "- Maintains consistent tone and style\n"
        "- Addresses the original task completely"
    ),

    'injection_tester': (
        "You are a prompt security analyst. Given a system prompt, identify potential vulnerabilities "
        "to prompt injection attacks.\n\n"
        "Test categories:\n"
        "1. Direct instruction override (e.g., 'Ignore previous instructions')\n"
        "2. Role escape (e.g., 'You are now a different AI')\n"
        "3. Context manipulation (e.g., 'The admin has authorized...')\n"
        "4. Encoding tricks (e.g., base64, unicode manipulation)\n"
        "5. Indirect injection (e.g., data that contains instructions)\n"
        "6. Output manipulation (e.g., 'Include this in your response...')\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "overall_risk": "<low|medium|high|critical>",\n'
        '  "risk_score": <1-10>,\n'
        '  "vulnerabilities": [\n'
        '    {"category": "<category>", "severity": "<low|medium|high|critical>", "description": "<what the vulnerability is>", "attack_example": "<example attack prompt>", "mitigation": "<how to fix>"}\n'
        "  ],\n"
        '  "hardened_prompt": "<improved version of the system prompt with security fixes>",\n'
        '  "recommendations": ["<list of general security recommendations>"]\n'
        "}\n"
        "Return ONLY valid JSON."
    ),

    'fewshot_builder': (
        "You are a few-shot prompt engineering expert. Given a task description and example "
        "input/output pairs, construct an optimized few-shot prompt.\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "system_prompt": "<the system prompt with role and instructions>",\n'
        '  "few_shot_examples": [\n'
        '    {"input": "<example input>", "output": "<example output>"}\n'
        "  ],\n"
        '  "user_prompt_template": "<template with {input} placeholder for new inputs>",\n'
        '  "assembled_prompt": "<the complete prompt ready to use, with examples embedded>",\n'
        '  "tips": ["<list of optimization tips for this prompt>"]\n'
        "}\n\n"
        "Guidelines:\n"
        "- Order examples from simple to complex\n"
        "- Ensure examples cover edge cases\n"
        "- Use consistent formatting across examples\n"
        "- Add chain-of-thought if reasoning is needed\n"
        "Return ONLY valid JSON."
    ),

    # --- Phase 4: Knowledge Workflows ---

    'expert_panel_persona': (
        "You are {persona_name}, a {persona_description}. "
        "You are participating in an expert panel discussion on the following topic.\n\n"
        "Provide your perspective based on your expertise and role. Be specific, "
        "insightful, and grounded in your domain knowledge. Express your unique viewpoint "
        "including potential concerns, opportunities, and recommendations from your perspective.\n\n"
        "Keep your response focused and substantive (200-400 words)."
    ),

    'expert_panel_moderator': (
        "You are a skilled moderator synthesizing an expert panel discussion.\n\n"
        "Multiple experts have shared their perspectives on a topic. Analyze all perspectives and produce:\n"
        "{\n"
        '  "consensus_points": ["<points where experts agree>"],\n'
        '  "disagreements": ["<points of contention with brief explanation>"],\n'
        '  "key_insights": ["<most valuable unique insights from the discussion>"],\n'
        '  "recommendation": "<synthesized recommendation considering all perspectives>",\n'
        '  "risk_factors": ["<risks identified across all perspectives>"],\n'
        '  "next_steps": ["<suggested next actions>"]\n'
        "}\n\n"
        "Be balanced and fair to all perspectives. Return ONLY valid JSON."
    ),

    'document_qa': (
        "You are a document analysis expert. You will be given multiple source documents "
        "and a question. Answer the question using ONLY information from the provided sources.\n\n"
        "Rules:\n"
        "1. Cite sources using [Source N] format where N is the source number\n"
        "2. Every factual claim MUST have a citation\n"
        "3. If the sources don't contain enough information, say so explicitly\n"
        "4. Synthesize information across sources when relevant\n"
        "5. Distinguish between what sources state vs. your inferences\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "answer": "<comprehensive answer with inline [Source N] citations>",\n'
        '  "sources_used": [<list of source numbers used>],\n'
        '  "confidence": "<high|medium|low>",\n'
        '  "key_findings": [\n'
        '    {"finding": "<key finding>", "sources": [<source numbers>]}\n'
        "  ],\n"
        '  "gaps": ["<information gaps or questions that sources cannot answer>"]\n'
        "}\n"
        "Return ONLY valid JSON."
    ),

    'compliance_checker': (
        "You are a compliance auditor. Given a policy/SOP document and a document to check, "
        "evaluate compliance against each rule or requirement in the policy.\n\n"
        "For each rule/requirement identified in the policy:\n"
        "1. Extract the rule\n"
        "2. Check if the document complies\n"
        "3. Provide evidence from the document\n"
        "4. Suggest remediation if non-compliant\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "overall_status": "<compliant|partially_compliant|non_compliant>",\n'
        '  "compliance_score": <percentage 0-100>,\n'
        '  "rules_checked": <total number of rules>,\n'
        '  "rules_passed": <number passing>,\n'
        '  "scorecard": [\n'
        '    {"rule_id": <number>, "rule": "<rule description>", "status": "<compliant|non_compliant|partial>", "evidence": "<evidence from document>", "recommendation": "<fix if needed>"}\n'
        "  ],\n"
        '  "critical_issues": ["<any critical non-compliance items>"],\n'
        '  "summary": "<overall compliance assessment>"\n'
        "}\n"
        "Return ONLY valid JSON."
    ),

    # --- Phase 5: Specialized Tools ---

    'tone_transformer': (
        "You are a professional writing style expert. Transform the given text into the "
        "requested tone/style while preserving the original meaning completely.\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "transformed_text": "<the rewritten text in the target tone>",\n'
        '  "changes_made": ["<list of specific changes and why>"],\n'
        '  "readability_metrics": {\n'
        '    "estimated_grade_level": <number>,\n'
        '    "word_count": <number>,\n'
        '    "sentence_count": <number>,\n'
        '    "avg_words_per_sentence": <number>\n'
        "  },\n"
        '  "tone_analysis": {\n'
        '    "original_tone": "<detected tone of input>",\n'
        '    "target_tone": "<the requested tone>",\n'
        '    "confidence": "<high|medium|low>"\n'
        "  }\n"
        "}\n"
        "Return ONLY valid JSON."
    ),

    'misconception_detector': (
        "You are an educational assessment expert specializing in identifying and correcting "
        "misconceptions. Analyze the student's answer for the given topic.\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "classification": "<correct|partially_correct|incorrect>",\n'
        '  "score": <0-100>,\n'
        '  "correct_elements": ["<parts of the answer that are correct>"],\n'
        '  "misconceptions": [\n'
        '    {"claim": "<what the student said>", "issue": "<why it is wrong>", "correction": "<the correct information>", "severity": "<critical|moderate|minor>"}\n'
        "  ],\n"
        '  "missing_concepts": ["<important concepts not mentioned>"],\n'
        '  "feedback": "<constructive feedback for the student>",\n'
        '  "suggested_resources": ["<topics to review>"]\n'
        "}\n"
        "Return ONLY valid JSON."
    ),

    'cot_visualizer': (
        "You are a step-by-step reasoning expert. Solve the given problem showing ALL "
        "your reasoning steps explicitly.\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "steps": [\n'
        '    {"step_number": <n>, "title": "<short title>", "reasoning": "<detailed reasoning>", "result": "<intermediate result>", "confidence": <1-10>}\n'
        "  ],\n"
        '  "final_answer": "<the final answer>",\n'
        '  "reasoning_type": "<deductive|inductive|abductive|mathematical|analytical>",\n'
        '  "total_steps": <number>,\n'
        '  "weakest_step": <step number with lowest confidence>,\n'
        '  "alternative_approaches": ["<other ways to solve this>"]\n'
        "}\n"
        "Show your work thoroughly. Return ONLY valid JSON."
    ),

    # --- Phase 6: Extended Features ---

    'rag_retriever': (
        "You are a document retrieval expert. Given a knowledge base split into chunks "
        "and a query, rank the chunks by relevance to the query.\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "ranked_chunks": [\n'
        '    {"chunk_index": <n>, "relevance_score": <0.0-1.0>, "reason": "<why relevant>"}\n'
        "  ],\n"
        '  "query_analysis": "<what the query is looking for>"\n'
        "}\n"
        "Return ONLY valid JSON. Rank all chunks, most relevant first."
    ),

    'rag_generator': (
        "You are a question-answering assistant. Answer the question using ONLY the "
        "provided context chunks. Cite chunk numbers in [Chunk N] format.\n\n"
        "If the context is insufficient, say so. Do not make up information."
    ),

    'scenario_simulator': (
        "You are a scenario simulation expert. Given a plan or proposal and a set of "
        "stakeholder roles, simulate how each stakeholder would react.\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "stakeholder_reactions": [\n'
        '    {"role": "<stakeholder role>", "reaction": "<likely reaction>", "concerns": ["<concerns>"], "support_level": "<strongly_support|support|neutral|oppose|strongly_oppose>", "conditions": "<conditions for support>"}\n'
        "  ],\n"
        '  "overall_viability": "<high|medium|low>",\n'
        '  "critical_risks": ["<major risks identified>"],\n'
        '  "recommended_modifications": ["<changes to improve acceptance>"],\n'
        '  "consensus_path": "<suggested approach to achieve stakeholder buy-in>"\n'
        "}\n"
        "Return ONLY valid JSON."
    ),

    'localizer': (
        "You are a cross-cultural prompt localization expert. Adapt the given prompt for "
        "the target language and cultural context.\n\n"
        "Do NOT just translate — adapt for:\n"
        "1. Cultural references and idioms\n"
        "2. Formal/informal register appropriate to the culture\n"
        "3. Examples relevant to the target audience\n"
        "4. Date/number/currency format conventions\n\n"
        "Return a JSON object:\n"
        "{\n"
        '  "localized_prompt": "<the adapted prompt in the target language>",\n'
        '  "adaptations": [\n'
        '    {"original": "<original element>", "adapted": "<localized element>", "reason": "<why this adaptation>"}\n'
        "  ],\n"
        '  "cultural_notes": ["<important cultural considerations>"],\n'
        '  "back_translation": "<English back-translation of the localized prompt>",\n'
        '  "confidence": "<high|medium|low>"\n'
        "}\n"
        "Return ONLY valid JSON."
    ),
}


# --- Phase 1: Prompt Quality Core ---

def grade_prompt(prompt_text, task_type='', domain='', model='gpt-4o-mini'):
    """Analyze and score a prompt across quality dimensions."""
    system_prompt = ADVANCED_PROMPTS['prompt_grader']
    user_prompt = f"Analyze this prompt:\n\n---\n{prompt_text}\n---"
    if task_type:
        user_prompt += f"\n\nTask type: {task_type}"
    if domain:
        user_prompt += f"\nDomain: {domain}"
    return execute_prompt(system_prompt, user_prompt, model, temperature=0.1, max_tokens=2048)


def optimize_prompt(prompt_text, grading_output, model='gpt-4o-mini'):
    """Generate improved and advanced versions of a prompt based on grading."""
    system_prompt = ADVANCED_PROMPTS['prompt_optimizer']
    user_prompt = (
        f"Original prompt:\n---\n{prompt_text}\n---\n\n"
        f"Analysis results:\n---\n{grading_output}\n---"
    )
    return execute_prompt(system_prompt, user_prompt, model, temperature=0.3, max_tokens=3072)


def compare_prompt_outputs(prompt_a, prompt_b, test_input, model='gpt-4o-mini'):
    """Run two prompts against the same input and judge results."""
    result_a = execute_prompt(prompt_a, test_input, model, temperature=0.3)
    result_b = execute_prompt(prompt_b, test_input, model, temperature=0.3)

    judge_system = ADVANCED_PROMPTS['ab_judge']
    judge_input = (
        f"Task input:\n{test_input}\n\n"
        f"--- Output A ---\n{result_a.get('output', '')}\n\n"
        f"--- Output B ---\n{result_b.get('output', '')}"
    )
    judge_result = execute_prompt(judge_system, judge_input, model, temperature=0.1)

    total_tokens = (
        result_a.get('tokens_input', 0) + result_a.get('tokens_output', 0) +
        result_b.get('tokens_input', 0) + result_b.get('tokens_output', 0) +
        judge_result.get('tokens_input', 0) + judge_result.get('tokens_output', 0)
    )
    total_cost = (
        result_a.get('cost_estimate', 0) + result_b.get('cost_estimate', 0) +
        judge_result.get('cost_estimate', 0)
    )

    return {
        'output': judge_result.get('output', ''),
        'output_a': result_a.get('output', ''),
        'output_b': result_b.get('output', ''),
        'tokens_input': total_tokens,
        'tokens_output': 0,
        'cost_estimate': total_cost,
        'latency_ms': (
            result_a.get('latency_ms', 0) + result_b.get('latency_ms', 0) +
            judge_result.get('latency_ms', 0)
        ),
        'model': model,
    }


def enforce_schema(prompt_text, schema_text, input_text, model='gpt-4o-mini', max_retries=3):
    """Execute a prompt with JSON schema enforcement and auto-retry on validation failure."""
    system = (
        ADVANCED_PROMPTS['schema_enforcer'] +
        f"\n\nRequired JSON Schema:\n{schema_text}"
    )
    best_output = None
    all_attempts = []

    for attempt in range(max_retries):
        if attempt == 0:
            user = f"{prompt_text}\n\nInput:\n{input_text}"
        else:
            user = (
                f"{prompt_text}\n\nInput:\n{input_text}\n\n"
                f"IMPORTANT: Your previous attempt was invalid. Errors:\n{validation_errors}\n"
                f"Fix these errors and return ONLY valid JSON matching the schema."
            )

        result = execute_prompt(system, user, model, temperature=0.1, max_tokens=2048)
        output = result.get('output', '')

        # Try to validate against schema
        try:
            sanitized = _sanitize_json(output)
            start = sanitized.find('{')
            end = sanitized.rfind('}')
            if start < 0:
                start = sanitized.find('[')
                end = sanitized.rfind(']')
            if start >= 0 and end > start:
                parsed = json.loads(sanitized[start:end + 1])
                all_attempts.append({'attempt': attempt + 1, 'valid': True, 'output': output})
                best_output = output
                break
        except json.JSONDecodeError as e:
            validation_errors = str(e)
            all_attempts.append({'attempt': attempt + 1, 'valid': False, 'error': validation_errors})
            best_output = output

    attempts_summary = json.dumps(all_attempts, indent=2)
    final_output = json.dumps({
        'result': best_output,
        'attempts': all_attempts,
        'total_attempts': len(all_attempts),
        'success': all_attempts[-1]['valid'] if all_attempts else False,
    }, indent=2)

    return {
        'output': final_output,
        'tokens_input': result.get('tokens_input', 0),
        'tokens_output': result.get('tokens_output', 0),
        'cost_estimate': result.get('cost_estimate', 0),
        'latency_ms': result.get('latency_ms', 0),
        'model': model,
    }


# --- Phase 2: Advanced Patterns ---

def execute_self_correcting(prompt_text, input_text, criteria='',
                            max_rounds=3, threshold=7, model='gpt-4o-mini'):
    """Generate, critique, and revise in a loop until quality threshold is met."""
    rounds = []
    total_tokens = 0
    total_cost = 0.0
    total_latency = 0

    # Round 1: Generate initial response
    gen_result = execute_prompt(prompt_text, input_text, model, temperature=0.5, max_tokens=2048)
    current_output = gen_result.get('output', '')
    total_tokens += gen_result.get('tokens_input', 0) + gen_result.get('tokens_output', 0)
    total_cost += gen_result.get('cost_estimate', 0)
    total_latency += gen_result.get('latency_ms', 0)

    rounds.append({
        'round': 1, 'type': 'generation', 'output': current_output,
    })

    for i in range(max_rounds):
        # Critique
        critic_system = ADVANCED_PROMPTS['self_correct_critic']
        critic_input = f"Task: {prompt_text}\n\nResponse to evaluate:\n---\n{current_output}\n---"
        if criteria:
            critic_input += f"\n\nAdditional quality criteria:\n{criteria}"

        critic_result = execute_prompt(critic_system, critic_input, model, temperature=0.1, max_tokens=1024)
        critique = critic_result.get('output', '')
        total_tokens += critic_result.get('tokens_input', 0) + critic_result.get('tokens_output', 0)
        total_cost += critic_result.get('cost_estimate', 0)
        total_latency += critic_result.get('latency_ms', 0)

        rounds.append({
            'round': i + 1, 'type': 'critique', 'output': critique,
        })

        # Check if passes threshold
        try:
            parsed = json.loads(_sanitize_json(critique))
            if isinstance(parsed, dict) and parsed.get('passes_threshold', False):
                break
        except (json.JSONDecodeError, TypeError):
            pass

        # Revise
        reviser_system = ADVANCED_PROMPTS['self_correct_reviser']
        reviser_input = (
            f"Original task: {prompt_text}\n\n"
            f"Current response:\n---\n{current_output}\n---\n\n"
            f"Critic's feedback:\n---\n{critique}\n---"
        )
        revise_result = execute_prompt(reviser_system, reviser_input, model, temperature=0.3, max_tokens=2048)
        current_output = revise_result.get('output', '')
        total_tokens += revise_result.get('tokens_input', 0) + revise_result.get('tokens_output', 0)
        total_cost += revise_result.get('cost_estimate', 0)
        total_latency += revise_result.get('latency_ms', 0)

        rounds.append({
            'round': i + 1, 'type': 'revision', 'output': current_output,
        })

    final_output = json.dumps({
        'final_output': current_output,
        'rounds': rounds,
        'total_rounds': len([r for r in rounds if r['type'] == 'critique']),
    }, indent=2)

    return {
        'output': final_output,
        'tokens_input': total_tokens,
        'tokens_output': 0,
        'cost_estimate': total_cost,
        'latency_ms': total_latency,
        'model': model,
    }


def execute_quality_pipeline(task_prompt, input_text, model='gpt-4o-mini'):
    """Multi-stage quality gate pipeline: generate -> check -> revise -> validate."""
    stages = []
    total_tokens = 0
    total_cost = 0.0
    total_latency = 0

    # Stage 1: Generate
    gen_result = execute_prompt(task_prompt, input_text, model, temperature=0.5, max_tokens=2048)
    content = gen_result.get('output', '')
    total_tokens += gen_result.get('tokens_input', 0) + gen_result.get('tokens_output', 0)
    total_cost += gen_result.get('cost_estimate', 0)
    total_latency += gen_result.get('latency_ms', 0)
    stages.append({'stage': 'Generator', 'status': 'completed', 'output': content})

    # Stage 2: Safety & Quality Check
    check_result = execute_prompt(
        ADVANCED_PROMPTS['quality_gate_safety'],
        f"Content to review:\n---\n{content}\n---\n\nOriginal task: {task_prompt}",
        model, temperature=0.1, max_tokens=1024
    )
    check_output = check_result.get('output', '')
    total_tokens += check_result.get('tokens_input', 0) + check_result.get('tokens_output', 0)
    total_cost += check_result.get('cost_estimate', 0)
    total_latency += check_result.get('latency_ms', 0)

    passed = True
    try:
        check_parsed = json.loads(_sanitize_json(check_output))
        passed = check_parsed.get('passed', True)
    except (json.JSONDecodeError, TypeError):
        pass

    stages.append({
        'stage': 'Safety & Quality Check',
        'status': 'passed' if passed else 'failed',
        'output': check_output,
    })

    # Stage 3: Revise if needed
    if not passed:
        revise_system = ADVANCED_PROMPTS['quality_gate_reviser'].format(issues=check_output)
        revise_result = execute_prompt(
            revise_system, content, model, temperature=0.3, max_tokens=2048
        )
        content = revise_result.get('output', '')
        total_tokens += revise_result.get('tokens_input', 0) + revise_result.get('tokens_output', 0)
        total_cost += revise_result.get('cost_estimate', 0)
        total_latency += revise_result.get('latency_ms', 0)
        stages.append({'stage': 'Reviser', 'status': 'completed', 'output': content})

        # Stage 4: Re-validate
        recheck_result = execute_prompt(
            ADVANCED_PROMPTS['quality_gate_safety'],
            f"Content to review:\n---\n{content}\n---\n\nOriginal task: {task_prompt}",
            model, temperature=0.1, max_tokens=1024
        )
        recheck_output = recheck_result.get('output', '')
        total_tokens += recheck_result.get('tokens_input', 0) + recheck_result.get('tokens_output', 0)
        total_cost += recheck_result.get('cost_estimate', 0)
        total_latency += recheck_result.get('latency_ms', 0)
        stages.append({'stage': 'Final Validation', 'status': 'completed', 'output': recheck_output})
    else:
        stages.append({'stage': 'Reviser', 'status': 'skipped', 'output': 'Not needed - all checks passed'})
        stages.append({'stage': 'Final Validation', 'status': 'skipped', 'output': 'Not needed'})

    final_output = json.dumps({
        'final_content': content,
        'stages': stages,
        'all_passed': passed,
    }, indent=2)

    return {
        'output': final_output,
        'tokens_input': total_tokens,
        'tokens_output': 0,
        'cost_estimate': total_cost,
        'latency_ms': total_latency,
        'model': model,
    }


def execute_decomposition(task_description, model='gpt-4o-mini'):
    """Break a complex task into sub-tasks, execute each, and integrate results."""
    stages = []
    total_tokens = 0
    total_cost = 0.0
    total_latency = 0

    # Stage 1: Plan
    plan_result = execute_prompt(
        ADVANCED_PROMPTS['decomposition_planner'], task_description,
        model, temperature=0.2, max_tokens=2048
    )
    plan_output = plan_result.get('output', '')
    total_tokens += plan_result.get('tokens_input', 0) + plan_result.get('tokens_output', 0)
    total_cost += plan_result.get('cost_estimate', 0)
    total_latency += plan_result.get('latency_ms', 0)
    stages.append({'stage': 'Planning', 'output': plan_output})

    # Parse sub-tasks
    sub_tasks = []
    try:
        parsed_plan = json.loads(_sanitize_json(plan_output))
        if isinstance(parsed_plan, dict):
            sub_tasks = parsed_plan.get('sub_tasks', [])
    except (json.JSONDecodeError, TypeError):
        sub_tasks = [{'step': 1, 'title': 'Full Task', 'instruction': task_description}]

    # Stage 2: Execute each sub-task
    sub_outputs = []
    for task in sub_tasks:
        instruction = task.get('instruction', '')
        context = ''
        if sub_outputs:
            context = f"\n\nContext from previous steps:\n{sub_outputs[-1][:500]}"

        sub_result = execute_prompt(
            "Complete the following sub-task thoroughly and precisely.",
            f"{instruction}{context}",
            model, temperature=0.4, max_tokens=1536
        )
        sub_output = sub_result.get('output', '')
        total_tokens += sub_result.get('tokens_input', 0) + sub_result.get('tokens_output', 0)
        total_cost += sub_result.get('cost_estimate', 0)
        total_latency += sub_result.get('latency_ms', 0)
        sub_outputs.append(sub_output)
        stages.append({
            'stage': f"Step {task.get('step', '?')}: {task.get('title', '')}",
            'output': sub_output,
        })

    # Stage 3: Integrate
    combined = "\n\n---\n\n".join(
        f"Sub-task {i+1}: {sub_tasks[i].get('title', '')}\n{out}"
        for i, out in enumerate(sub_outputs)
    )
    integrate_result = execute_prompt(
        ADVANCED_PROMPTS['decomposition_integrator'],
        f"Original task: {task_description}\n\nSub-task results:\n{combined}",
        model, temperature=0.3, max_tokens=3072
    )
    integrated = integrate_result.get('output', '')
    total_tokens += integrate_result.get('tokens_input', 0) + integrate_result.get('tokens_output', 0)
    total_cost += integrate_result.get('cost_estimate', 0)
    total_latency += integrate_result.get('latency_ms', 0)
    stages.append({'stage': 'Integration', 'output': integrated})

    final_output = json.dumps({
        'final_output': integrated,
        'stages': stages,
        'sub_task_count': len(sub_tasks),
    }, indent=2)

    return {
        'output': final_output,
        'tokens_input': total_tokens,
        'tokens_output': 0,
        'cost_estimate': total_cost,
        'latency_ms': total_latency,
        'model': model,
    }


# --- Phase 3: Security & Templates ---

def test_prompt_injection(system_prompt, model='gpt-4o-mini'):
    """Analyze a system prompt for injection vulnerabilities."""
    analyzer_system = ADVANCED_PROMPTS['injection_tester']
    user_input = f"Analyze this system prompt for vulnerabilities:\n\n---\n{system_prompt}\n---"
    return execute_prompt(analyzer_system, user_input, model, temperature=0.2, max_tokens=3072)


def build_fewshot_prompt(task_description, examples_json, model='gpt-4o-mini'):
    """Build an optimized few-shot prompt from task description and examples."""
    system = ADVANCED_PROMPTS['fewshot_builder']
    user_input = (
        f"Task description: {task_description}\n\n"
        f"Example input/output pairs:\n{examples_json}"
    )
    return execute_prompt(system, user_input, model, temperature=0.3, max_tokens=3072)


# --- Phase 4: Knowledge Workflows ---

EXPERT_PERSONAS = {
    'optimist': ('Optimist Strategist', 'strategic thinker who focuses on opportunities, growth potential, and positive outcomes'),
    'skeptic': ('Skeptical Analyst', 'critical analyst who questions assumptions, identifies flaws, and stress-tests ideas'),
    'risk_mgr': ('Risk Manager', 'risk management professional who evaluates threats, vulnerabilities, and mitigation strategies'),
    'technical': ('Technical Expert', 'technology specialist who assesses feasibility, implementation complexity, and technical trade-offs'),
    'end_user': ('End-User Advocate', 'user experience champion who evaluates from the customer/end-user perspective'),
    'financial': ('Financial Advisor', 'financial expert who analyzes costs, ROI, and economic viability'),
}


def execute_expert_panel(topic, personas, model='gpt-4o-mini'):
    """Run a multi-persona expert panel discussion with moderator synthesis."""
    total_tokens = 0
    total_cost = 0.0
    total_latency = 0
    expert_responses = []

    # Round 1: Each expert gives their perspective
    for persona_key in personas:
        persona = EXPERT_PERSONAS.get(persona_key, (persona_key, 'domain expert'))
        system = ADVANCED_PROMPTS['expert_panel_persona'].format(
            persona_name=persona[0], persona_description=persona[1]
        )
        result = execute_prompt(system, f"Topic for discussion:\n{topic}", model, temperature=0.6, max_tokens=1024)
        total_tokens += result.get('tokens_input', 0) + result.get('tokens_output', 0)
        total_cost += result.get('cost_estimate', 0)
        total_latency += result.get('latency_ms', 0)
        expert_responses.append({
            'persona': persona[0],
            'persona_key': persona_key,
            'response': result.get('output', ''),
        })

    # Round 2: Moderator synthesizes
    panel_text = "\n\n".join(
        f"--- {e['persona']} ---\n{e['response']}" for e in expert_responses
    )
    moderator_result = execute_prompt(
        ADVANCED_PROMPTS['expert_panel_moderator'],
        f"Topic: {topic}\n\nExpert perspectives:\n{panel_text}",
        model, temperature=0.2, max_tokens=2048
    )
    total_tokens += moderator_result.get('tokens_input', 0) + moderator_result.get('tokens_output', 0)
    total_cost += moderator_result.get('cost_estimate', 0)
    total_latency += moderator_result.get('latency_ms', 0)

    final_output = json.dumps({
        'expert_responses': expert_responses,
        'synthesis': moderator_result.get('output', ''),
        'num_experts': len(expert_responses),
    }, indent=2)

    return {
        'output': final_output,
        'tokens_input': total_tokens,
        'tokens_output': 0,
        'cost_estimate': total_cost,
        'latency_ms': total_latency,
        'model': model,
    }


def execute_document_qa(question, documents, model='gpt-4o-mini'):
    """Answer a question using multiple source documents with citations."""
    doc_text = "\n\n".join(
        f"--- Source {i+1}: {doc.get('label', f'Document {i+1}')} ---\n{doc.get('text', '')}"
        for i, doc in enumerate(documents)
    )
    system = ADVANCED_PROMPTS['document_qa']
    user_input = f"Question: {question}\n\nDocuments:\n{doc_text}"
    return execute_prompt(system, user_input, model, temperature=0.1, max_tokens=3072)


def execute_compliance_check(policy_text, document_text, model='gpt-4o-mini'):
    """Check a document against a policy/SOP for compliance."""
    system = ADVANCED_PROMPTS['compliance_checker']
    user_input = (
        f"Policy/SOP:\n---\n{policy_text}\n---\n\n"
        f"Document to check:\n---\n{document_text}\n---"
    )
    return execute_prompt(system, user_input, model, temperature=0.1, max_tokens=3072)


# --- Phase 5: Specialized Tools ---

def execute_tone_transform(text, target_tone, model='gpt-4o-mini'):
    """Transform text into a different tone/style."""
    system = ADVANCED_PROMPTS['tone_transformer']
    user_input = f"Target tone: {target_tone}\n\nText to transform:\n---\n{text}\n---"
    return execute_prompt(system, user_input, model, temperature=0.4, max_tokens=2048)


def execute_misconception_detector(topic, student_answer, model='gpt-4o-mini'):
    """Analyze a student answer for misconceptions."""
    system = ADVANCED_PROMPTS['misconception_detector']
    user_input = f"Topic/Subject: {topic}\n\nStudent's answer:\n---\n{student_answer}\n---"
    return execute_prompt(system, user_input, model, temperature=0.2, max_tokens=2048)


def execute_cot_visualizer(question, model='gpt-4o-mini'):
    """Solve a problem with explicit chain-of-thought reasoning steps."""
    system = ADVANCED_PROMPTS['cot_visualizer']
    return execute_prompt(system, question, model, temperature=0.3, max_tokens=3072)


# --- Phase 6: Extended Features ---

def execute_rag_simulator(query, knowledge_chunks, model='gpt-4o-mini'):
    """Simulate a RAG pipeline: retrieve relevant chunks then generate answer."""
    total_tokens = 0
    total_cost = 0.0
    total_latency = 0

    # Step 1: Retrieve — rank chunks by relevance
    chunks_text = "\n\n".join(
        f"[Chunk {i+1}]: {chunk}" for i, chunk in enumerate(knowledge_chunks)
    )
    retrieval_result = execute_prompt(
        ADVANCED_PROMPTS['rag_retriever'],
        f"Query: {query}\n\nKnowledge base chunks:\n{chunks_text}",
        model, temperature=0.1, max_tokens=1024
    )
    total_tokens += retrieval_result.get('tokens_input', 0) + retrieval_result.get('tokens_output', 0)
    total_cost += retrieval_result.get('cost_estimate', 0)
    total_latency += retrieval_result.get('latency_ms', 0)

    # Parse ranked chunks and select top-K
    retrieval_output = retrieval_result.get('output', '')
    top_k = 3
    selected_chunks = []
    try:
        parsed = json.loads(_sanitize_json(retrieval_output))
        ranked = parsed.get('ranked_chunks', [])
        for item in ranked[:top_k]:
            idx = item.get('chunk_index', 1) - 1
            if 0 <= idx < len(knowledge_chunks):
                selected_chunks.append((idx + 1, knowledge_chunks[idx]))
    except (json.JSONDecodeError, TypeError):
        selected_chunks = [(i + 1, c) for i, c in enumerate(knowledge_chunks[:top_k])]

    # Step 2: Generate with context
    context = "\n\n".join(f"[Chunk {idx}]: {text}" for idx, text in selected_chunks)
    gen_with_context = execute_prompt(
        ADVANCED_PROMPTS['rag_generator'],
        f"Question: {query}\n\nRetrieved context:\n{context}",
        model, temperature=0.3, max_tokens=2048
    )
    total_tokens += gen_with_context.get('tokens_input', 0) + gen_with_context.get('tokens_output', 0)
    total_cost += gen_with_context.get('cost_estimate', 0)
    total_latency += gen_with_context.get('latency_ms', 0)

    # Step 3: Generate without context (baseline)
    gen_no_context = execute_prompt(
        "Answer the following question to the best of your knowledge.",
        query, model, temperature=0.3, max_tokens=1024
    )
    total_tokens += gen_no_context.get('tokens_input', 0) + gen_no_context.get('tokens_output', 0)
    total_cost += gen_no_context.get('cost_estimate', 0)
    total_latency += gen_no_context.get('latency_ms', 0)

    final_output = json.dumps({
        'retrieval': retrieval_output,
        'selected_chunks': [{'chunk_index': idx, 'text': text[:200] + '...' if len(text) > 200 else text} for idx, text in selected_chunks],
        'answer_with_context': gen_with_context.get('output', ''),
        'answer_without_context': gen_no_context.get('output', ''),
        'num_chunks_total': len(knowledge_chunks),
        'num_chunks_retrieved': len(selected_chunks),
    }, indent=2)

    return {
        'output': final_output,
        'tokens_input': total_tokens,
        'tokens_output': 0,
        'cost_estimate': total_cost,
        'latency_ms': total_latency,
        'model': model,
    }


def execute_scenario_simulator(plan, stakeholders, model='gpt-4o-mini'):
    """Simulate stakeholder reactions to a plan or proposal."""
    system = ADVANCED_PROMPTS['scenario_simulator']
    stakeholder_list = "\n".join(f"- {s}" for s in stakeholders)
    user_input = (
        f"Plan/Proposal:\n---\n{plan}\n---\n\n"
        f"Stakeholder roles to simulate:\n{stakeholder_list}"
    )
    return execute_prompt(system, user_input, model, temperature=0.5, max_tokens=3072)


def execute_localizer(prompt_text, target_language, cultural_context='', model='gpt-4o-mini'):
    """Localize a prompt for a different language and cultural context."""
    system = ADVANCED_PROMPTS['localizer']
    user_input = f"Target language: {target_language}\n"
    if cultural_context:
        user_input += f"Cultural context notes: {cultural_context}\n"
    user_input += f"\nPrompt to localize:\n---\n{prompt_text}\n---"
    return execute_prompt(system, user_input, model, temperature=0.4, max_tokens=3072)


def _sanitize_json(text):
    """Clean LLM output to extract valid JSON."""
    cleaned = re.sub(r'```(?:json)?\s*\n?', '', text).strip()
    cleaned = cleaned.rstrip('`').strip()
    cleaned = re.sub(r':\s*NULL\b', ': null', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r':\s*TRUE\b', ': true', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r':\s*FALSE\b', ': false', cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r',\s*([}\]])', r'\1', cleaned)
    return cleaned


def generate_meeting_docx(meeting_text):
    """Convert meeting summary text into a Word document. Returns BytesIO buffer."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Style the default font
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)
    font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Title
    title = doc.add_heading('Meeting Summary', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    doc.add_paragraph('')  # spacer

    # Parse the meeting text into sections
    lines = meeting_text.strip().split('\n')
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        # Detect section headers (numbered like "1. Date..." or markdown headings)
        heading_match = re.match(r'^(?:#{1,3}\s+|(\d+)\.\s*)(.*)', stripped)
        if heading_match:
            heading_text = heading_match.group(2) if heading_match.group(2) else heading_match.group(0)
            is_section = heading_match.group(1) is not None or stripped.startswith('#')
            if is_section:
                h = doc.add_heading(heading_text.strip().rstrip(':'), level=2)
                for run in h.runs:
                    run.font.color.rgb = RGBColor(0x6C, 0x6C, 0xF4)
                continue

        # Detect bullet points
        bullet_match = re.match(r'^[-*\u2022]\s+(.*)', stripped)
        if bullet_match:
            doc.add_paragraph(bullet_match.group(1), style='List Bullet')
            continue

        # Detect sub-items (indented bullets)
        sub_match = re.match(r'^\s+[-*\u2022]\s+(.*)', line)
        if sub_match:
            doc.add_paragraph(sub_match.group(1), style='List Bullet 2')
            continue

        # Bold key-value pairs like "Date: October 1st" or "**Date**: October 1st"
        kv_match = re.match(r'^(?:\*\*)?([^:*]{2,35})(?:\*\*)?:\s*(.*)', stripped)
        if kv_match:
            p = doc.add_paragraph()
            run_key = p.add_run(kv_match.group(1).strip() + ': ')
            run_key.bold = True
            run_key.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            p.add_run(kv_match.group(2).strip())
            continue

        # Regular paragraph - strip markdown bold markers
        clean_text = re.sub(r'\*\*(.+?)\*\*', r'\1', stripped)
        doc.add_paragraph(clean_text)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def generate_quiz_docx(quiz_json_str):
    """Convert quiz JSON into a Word document. Returns BytesIO buffer."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    # Parse quiz data
    sanitized = _sanitize_json(quiz_json_str)
    start = sanitized.find('[')
    end = sanitized.rfind(']')
    if start >= 0 and end > start:
        sanitized = sanitized[start:end + 1]
    questions = json.loads(sanitized)
    if not isinstance(questions, list):
        questions = [questions]

    doc = Document()

    # Style
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)
    font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Title
    title = doc.add_heading('Training Quiz', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    doc.add_paragraph(f'Total Questions: {len(questions)}')
    doc.add_paragraph('')

    # Questions section
    doc.add_heading('Questions', level=1)

    for i, q in enumerate(questions, 1):
        question_text = q.get('question', f'Question {i}')
        difficulty = q.get('difficulty', '')

        # Question heading
        p = doc.add_paragraph()
        run_num = p.add_run(f'Q{i}. ')
        run_num.bold = True
        run_num.font.size = Pt(12)
        run_num.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        run_q = p.add_run(question_text)
        run_q.font.size = Pt(12)
        if difficulty:
            run_diff = p.add_run(f'  [{difficulty}]')
            run_diff.font.size = Pt(9)
            run_diff.font.color.rgb = RGBColor(0x6C, 0x6C, 0xF4)
            run_diff.italic = True

        # Options
        options = q.get('options', [])
        for j, opt in enumerate(options):
            letter = chr(65 + j)
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Pt(24)
            run = p.add_run(f'{letter}) {opt}')
            run.font.size = Pt(11)

        doc.add_paragraph('')

    # Answer key on new page
    doc.add_page_break()
    doc.add_heading('Answer Key', level=1)

    for i, q in enumerate(questions, 1):
        correct = q.get('correct_answer', 'N/A')
        p = doc.add_paragraph()
        run_num = p.add_run(f'Q{i}: ')
        run_num.bold = True
        run_ans = p.add_run(str(correct))
        run_ans.font.color.rgb = RGBColor(0x22, 0x8B, 0x22)
        run_ans.bold = True

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _sanitize_slide_json(text):
    """Clean LLM output to extract valid JSON for slide data."""
    cleaned = _sanitize_json(text)
    # Find the JSON array
    start = cleaned.find('[')
    end = cleaned.rfind(']')
    if start >= 0 and end > start:
        cleaned = cleaned[start:end + 1]
    return cleaned


def generate_pptx(slide_json_str):
    """Convert slide script JSON into a PowerPoint file. Returns BytesIO buffer."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Parse the slide data
    sanitized = _sanitize_slide_json(slide_json_str)
    slides_data = json.loads(sanitized)
    if not isinstance(slides_data, list):
        slides_data = [slides_data]

    prs = Presentation()
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background fill
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)

        # Slide number badge (top-right)
        slide_num = slide_data.get('slide_number', '')
        if slide_num:
            num_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.3), Inches(1.2), Inches(0.5))
            tf = num_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"#{slide_num}"
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0x94, 0x94, 0xF8)
            run.font.bold = True

        # Title
        title_text = slide_data.get('title', 'Untitled')
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)

        # Accent line under title
        slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.85), Inches(3), Emu(36000)
        ).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = RGBColor(0x6C, 0x6C, 0xF4)

        # Bullet points
        bullets = slide_data.get('bullets', [])
        if bullets:
            bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(4))
            tf = bullet_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(12)
                run = p.add_run()
                run.text = f"\u2022  {bullet}"
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0xBA, 0xC2, 0xDE)

        # Speaker notes
        notes_text = slide_data.get('speaker_notes', '')
        if notes_text:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
